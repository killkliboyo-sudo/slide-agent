from __future__ import annotations

"""Assembler renders a minimal PPTX and companion markdown preview."""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .models import AssemblyResult, AssetSpec, SlideDraft


def assemble(slides: list[SlideDraft], requested_output: Path, theme: Optional[dict[str, str]] = None) -> AssemblyResult:
    """Render the slide drafts into a PPTX and a lightweight markdown preview."""
    requested_output = requested_output.expanduser()
    requested_output.parent.mkdir(parents=True, exist_ok=True)

    presentation = _build_presentation(slides, theme or {})
    presentation.save(requested_output)

    preview_path = requested_output.with_suffix(".md")
    preview_path.write_text(_render_preview(slides), encoding="utf-8")

    notes = (
        "Generated PPTX with simple title/body layouts; markdown preview emitted "
        "for quick inspection."
    )
    return AssemblyResult(
        requested_output=requested_output,
        preview_path=preview_path,
        slides_built=len(slides),
        notes=notes,
    )


def _build_presentation(slides: list[SlideDraft], theme: dict[str, str]) -> Presentation:
    """Create a presentation with 16:9 ratio and basic styling."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    font_name = theme.get("font", "Segoe UI")
    palette = _resolve_palette(theme.get("palette"))

    for slide in slides:
        _add_slide(prs, slide, font_name, palette)
    return prs


def _add_slide(prs: Presentation, slide: SlideDraft, font_name: str, palette: dict[str, str]) -> None:
    """Add a single slide honoring the requested layout when possible."""
    layout = prs.slide_layouts[6]  # blank layout
    ppt_slide = prs.slides.add_slide(layout)

    _apply_background(ppt_slide, palette)

    margin = Inches(0.6)
    body_top = margin + Inches(1.0)
    body_height = prs.slide_height - body_top - Inches(1.2)
    body_width = prs.slide_width - (margin * 2)
    gap = Inches(0.4)

    _add_title(ppt_slide, slide.title, margin, margin, body_width, font_name, palette["accent"])

    if slide.layout == "split":
        text_width = (body_width - gap) * 0.55
        _add_bullets_box(
            ppt_slide,
            slide.bullets,
            left=margin,
            top=body_top,
            width=text_width,
            height=body_height,
            font_name=font_name,
            color=palette["primary"],
        )
        _add_asset_placeholder(
            ppt_slide,
            slide.assets,
            left=margin + text_width + gap,
            top=body_top,
            width=body_width - text_width - gap,
            height=body_height,
            font_name=font_name,
            color=palette["muted"],
        )
    elif slide.layout == "stacked":
        _add_bullets_box(
            ppt_slide,
            slide.bullets,
            left=margin,
            top=body_top,
            width=body_width,
            height=body_height * 0.55,
            font_name=font_name,
            color=palette["primary"],
        )
        _add_asset_placeholder(
            ppt_slide,
            slide.assets,
            left=margin,
            top=body_top + body_height * 0.6,
            width=body_width,
            height=body_height * 0.35,
            font_name=font_name,
            color=palette["muted"],
        )
    else:  # focus
        _add_bullets_box(
            ppt_slide,
            slide.bullets,
            left=margin,
            top=body_top + Inches(0.4),
            width=body_width,
            height=body_height * 0.8,
            emphasize=True,
            font_name=font_name,
            color=palette["primary"],
        )

    _add_footer(ppt_slide, slide.sources, slide.notes or "", font_name, color=palette["muted"])


def _add_title(slide, title: str, left, top, width, font_name: str, color: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    frame = box.text_frame
    frame.text = title
    para = frame.paragraphs[0]
    para.font.size = Pt(32)
    para.font.bold = True
    para.font.name = font_name
    para.font.color.rgb = _hex_to_rgb(color)


def _add_bullets_box(
    slide,
    bullets: list[str],
    left,
    top,
    width,
    height,
    emphasize: bool = False,
    font_name: str = "Segoe UI",
    color: str = "#000000",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets or ["Highlight key message"]):
        paragraph = tf.add_paragraph()
        if idx == 0:
            paragraph.level = 0
        paragraph.text = bullet
        paragraph.font.size = Pt(24 if emphasize else 22)
        paragraph.font.name = font_name
        paragraph.font.color.rgb = _hex_to_rgb(color)
    if emphasize:
        for paragraph in tf.paragraphs:
            paragraph.font.bold = True


def _add_asset_placeholder(
    slide, assets: list[AssetSpec], left, top, width, height, font_name: str, color: str = "#666666"
) -> None:
    if not assets:
        return
    asset = assets[0]
    if asset.path and asset.path.exists():
        try:
            slide.shapes.add_picture(str(asset.path), left, top, width=width, height=height)
            return
        except Exception:
            pass

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    description = asset.prompt or f"{asset.type} placeholder"
    frame.text = description
    frame.paragraphs[0].font.size = Pt(18)
    frame.paragraphs[0].font.name = font_name
    frame.paragraphs[0].font.italic = True
    frame.paragraphs[0].font.color.rgb = _hex_to_rgb(color)


def _add_footer(slide, sources: list[str], note: str, font_name: str, color: str) -> None:
    footer = slide.shapes.add_textbox(
        Inches(0.6),
        slide.part.slide_height - Inches(0.8),
        slide.part.slide_width - Inches(1.2),
        Inches(0.6),
    )
    frame = footer.text_frame
    frame.clear()
    if sources:
        src_para = frame.add_paragraph()
        src_para.text = f"Sources: {', '.join(sources)}"
        src_para.font.size = Pt(12)
        src_para.font.name = font_name
        src_para.font.color.rgb = _hex_to_rgb(color)
    if note:
        note_para = frame.add_paragraph()
        note_para.text = note
        note_para.font.size = Pt(12)
        note_para.font.name = font_name
        note_para.font.color.rgb = _hex_to_rgb(color)


def _apply_background(slide, palette: dict[str, str]) -> None:
    """Apply a solid background and a soft overlay strip for contrast."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(palette["background"])

    # Add a translucent accent bar at the top for the title.
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slide.part.slide_width,
        Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(palette["bar"])
    bar.fill.fore_color.brightness = -0.15
    bar.line.fill.background()


def _render_preview(slides: list[SlideDraft]) -> str:
    lines = ["# Auto Presentation Agent Preview", ""]
    for idx, slide in enumerate(slides, start=1):
        lines.append(f"## Slide {idx}: {slide.title}")
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        if slide.assets:
            assets = ", ".join(asset.type for asset in slide.assets)
            lines.append(f"_Assets_: {assets}")
        if slide.sources:
            lines.append(f"_Sources_: {', '.join(slide.sources)}")
        if slide.notes:
            lines.append(f"_Notes_: {slide.notes}")
        lines.append("")
    return "\n".join(lines)


def _hex_to_rgb(value: str):
    """Convert #rrggbb to pptx RGBColor."""
    value = value.lstrip("#")
    r, g, b = int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
    from pptx.dml.color import RGBColor

    return RGBColor(r, g, b)


def _resolve_palette(name: Optional[str]) -> dict[str, str]:
    """Return a palette dict keyed by semantic color roles."""
    base = {
        "background": "#0f172a",  # slate-900
        "bar": "#0b1221",
        "accent": "#eab308",  # amber-500
        "primary": "#e2e8f0",  # slate-200
        "muted": "#94a3b8",  # slate-400
    }
    if name and name.lower() == "light":
        return {
            "background": "#f8fafc",
            "bar": "#e2e8f0",
            "accent": "#0f172a",
            "primary": "#0f172a",
            "muted": "#475569",
        }
    return base
