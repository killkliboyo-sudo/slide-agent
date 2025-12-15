from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

import types

from auto_presentation_agent.designer import design_slides
from auto_presentation_agent import data_analysis as da
from auto_presentation_agent import imagegen
from auto_presentation_agent import llm
from auto_presentation_agent.models import (
    AssetSpec,
    ContentSummary,
    OutlinePlan,
    OutlineSlide,
    PresentationRequest,
    SlideDraft,
)
from auto_presentation_agent.outline import generate_outline
try:
    from auto_presentation_agent.pipeline import run_pipeline

    PPTX_AVAILABLE = True
except ModuleNotFoundError:
    run_pipeline = None
    PPTX_AVAILABLE = False


class OutlineDesignerTests(unittest.TestCase):
    def test_outline_respects_duration_and_visual_choice(self) -> None:
        summary = ContentSummary(
            topics=["Alpha", "Beta"],
            findings=["Result highlights", "Secondary observation"],
            data_points=["metric_a mean 1.2"],
        )
        outline = generate_outline(summary, duration_minutes=2, style_prefs={"palette": "dark"})

        self.assertEqual(len(outline.slides), 2)
        self.assertTrue(all(slide.title.endswith("key takeaway") for slide in outline.slides))
        self.assertTrue(all(slide.visual_suggestion == "chart" for slide in outline.slides))
        self.assertEqual(outline.theme.get("palette"), "dark")
        self.assertEqual(outline.theme.get("font"), "Segoe UI")

    def test_designer_cycles_layouts_and_converts_titles(self) -> None:
        outline = OutlinePlan(
            slides=[
                OutlineSlide(title="Alpha: conclusion here", bullets=["one", "two"], visual_suggestion="image"),
                OutlineSlide(title="Beta insight", bullets=["only bullet"], visual_suggestion=None),
                OutlineSlide(title="Gamma: findings", bullets=["x"], visual_suggestion="chart"),
            ]
        )
        drafts = design_slides(outline)

        self.assertEqual([d.layout for d in drafts], ["split", "stacked", "focus"])
        self.assertEqual(drafts[0].title, "conclusion here")
        self.assertEqual(drafts[1].title, "Beta insight")
        self.assertEqual(drafts[0].assets[0].type, "image")
        self.assertEqual(drafts[2].assets[0].type, "chart")

    def test_designer_can_attach_generated_image(self) -> None:
        outline = OutlinePlan(
            slides=[
                OutlineSlide(title="Alpha: conclusion here", bullets=["one"], visual_suggestion="image"),
            ]
        )
        tmp = Path.cwd() / "assets_test"
        if tmp.exists():
            for f in tmp.glob("*"):
                f.unlink()
        tmp.mkdir(exist_ok=True)

        def fake_gen(prompt, assets_dir):
            p = assets_dir / "img.png"
            p.write_bytes(b"\x89PNG\r\n\x1a\n")
            return p

        drafts = design_slides(outline, fake_gen, tmp)
        self.assertTrue(drafts[0].assets[0].path and drafts[0].assets[0].path.exists())
        for f in tmp.glob("*"):
            f.unlink()
        tmp.rmdir()

@unittest.skipUnless(PPTX_AVAILABLE, "python-pptx is required for the pipeline smoke test")
class PipelineSmokeTests(unittest.TestCase):
    def test_end_to_end_pipeline_creates_outputs(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            base = Path(tmpdir)
            text_file = base / "note.txt"
            text_file.write_text("This is a short note about quarterly progress.", encoding="utf-8")
            csv_file = base / "table.csv"
            csv_file.write_text("value\n1\n2\n3\n4\n", encoding="utf-8")
            output_path = base / "deck.pptx"

            preferred_font = "Inter"
            palette_name = "light"
            request = PresentationRequest(
                inputs=[text_file, csv_file],
                instructions="Focus on results and upcoming work.",
                duration_minutes=2,
                style_prefs={"palette": palette_name, "font": preferred_font},
                output_path=output_path,
            )

            result = run_pipeline(request)

            self.assertTrue(output_path.exists())
            preview_path = output_path.with_suffix(".md")
            self.assertTrue(preview_path.exists())
            self.assertGreaterEqual(result.slides_built, 1)
            preview_text = preview_path.read_text(encoding="utf-8")
            self.assertIn("Slide 1", preview_text)

            from pptx import Presentation as PptPresentation

            prs = PptPresentation(output_path)
            fonts = []
            backgrounds = []
            for shape in prs.slides[0].shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for para in shape.text_frame.paragraphs:
                    if para.font and para.font.name:
                        fonts.append(para.font.name)
            # collect background color of first slide
            bg_fill = prs.slides[0].background.fill
            if bg_fill.type == 1:  # solid fill
                color = bg_fill.fore_color.rgb
                backgrounds.append((color[0], color[1], color[2]))
            self.assertIn(preferred_font, fonts)
            self.assertTrue(backgrounds)  # palette applied to background


if __name__ == "__main__":
    unittest.main()
class DataAnalysisPdfTests(unittest.TestCase):
    def test_pdf_is_noted_and_summarized_or_warned(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            base = Path(tmpdir)
            pdf_path = base / "sample.pdf"
            # Minimal placeholder content; parsing may warn if pdfplumber is available.
            pdf_path.write_text("%PDF-1.4\n1 0 obj\n<<>>\nendobj\n", encoding="utf-8")

            request = PresentationRequest(inputs=[pdf_path])
            summary = da.analyze_request(request)

            self.assertIn(pdf_path.stem, summary.topics)
            findings_text = " ".join(summary.findings).lower()
            self.assertIn("pdf", findings_text)


class LlmIntegrationTests(unittest.TestCase):
    def test_llm_summary_is_appended_when_enabled(self) -> None:
        fake_llm = types.SimpleNamespace(generate=lambda prompt: "core insights summarized")
        request = PresentationRequest(inputs=[], use_llm=True)
        summary = da.analyze_request(request, llm_client=fake_llm)
        joined = " ".join(summary.findings).lower()
        self.assertIn("llm summary", joined)

    def test_gemini_image_backend_uses_client_bytes(self) -> None:
        tmp = Path(tempfile.mkdtemp())
        try:
            prompt = "abstract waves"
            fake_client = types.SimpleNamespace(
                generate=lambda p: "text",
                generate_image=lambda p, image_model=None: b"\x89PNG\r\n\x1a\n",
            )
            path = imagegen.generate_image(prompt, tmp, backend="gemini", gemini_client=fake_client)
            self.assertTrue(path and path.exists())
        finally:
            for f in tmp.glob("*"):
                f.unlink()
            tmp.rmdir()

    def test_gemini_model_listing_parses_names(self) -> None:
        fake_payload = {"models": [{"name": "models/gemini-1"}, {"name": "models/gemini-2"}]}

        def fake_urlopen(req, timeout=10):
            class _Resp:
                status = 200

                def __enter__(self):
                    return self

                def __exit__(self, exc_type, exc, tb):
                    return False

                def read(self):
                    import json

                    return json.dumps(fake_payload).encode("utf-8")

            return _Resp()

        with mock.patch("urllib.request.urlopen", fake_urlopen):
            models = llm.list_gemini_models("dummy")
        self.assertEqual(models, ["models/gemini-1", "models/gemini-2"])


class PreviewSnapshotTests(unittest.TestCase):
    def test_markdown_preview_matches_golden(self) -> None:
        from auto_presentation_agent.assembler import _render_preview

        slides = [
            SlideDraft(
                title="Example Title",
                bullets=["First point", "Second point"],
                layout="split",
                assets=[AssetSpec(type="image", prompt="something")],
                notes="Drafted per SPEC guidelines.",
            )
        ]
        preview = _render_preview(slides).strip()
        golden = (Path(__file__).parent / "fixtures" / "preview_golden.md").read_text(encoding="utf-8").strip()
        self.assertEqual(preview, golden)
